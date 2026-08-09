import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Constants
    C_DARK_BG = RGBColor(0, 45, 28)        # #002d1c
    C_PINE = RGBColor(7, 71, 46)          # #07472e
    C_EMERALD = RGBColor(27, 67, 50)       # #1b4332
    C_MUTED_GREEN = RGBColor(62, 103, 82)  # #3e6752
    C_CANVAS = RGBColor(252, 249, 248)    # #fcf9f8
    C_WHITE = RGBColor(255, 255, 255)
    C_TEXT = RGBColor(27, 28, 28)          # #1b1c1c
    C_TEXT_MUTED = RGBColor(113, 121, 115) # #717973
    C_BORDER = RGBColor(229, 226, 217)    # #e5e2d9
    C_BLUE_ACCENT = RGBColor(14, 165, 233) # #0ea5e9
    C_EMERALD_ACCENT = RGBColor(16, 185, 129) # #10b981
    
    # Priority Colors
    C_VERY_HIGH = RGBColor(186, 26, 26)   # #ba1a1a
    C_HIGH = RGBColor(244, 140, 36)       # #f48c24
    C_MEDIUM = RGBColor(234, 181, 82)     # #eab552
    C_LOW = RGBColor(127, 163, 143)       # #7fa38f

    def set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape_card(slide, left, top, width, height, bg_color, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def add_header(slide, title, subtitle=None, dark=False):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.name = 'Inter'
        p.font.color.rgb = C_WHITE if dark else C_TEXT
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(14)
            p2.font.name = 'Inter'
            p2.font.color.rgb = RGBColor(180, 210, 195) if dark else C_TEXT_MUTED
            p2.space_before = Pt(4)

    def add_footer(slide, connecting_text, slide_num, dark=False):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{connecting_text}"'
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.name = 'Inter'
        p.font.color.rgb = RGBColor(160, 200, 180) if dark else C_TEXT_MUTED
        
        # Slide counter right-aligned
        txBox2 = slide.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.0), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f'{slide_num:02d} / 14'
        p2.alignment = PP_ALIGN.RIGHT
        p2.font.size = Pt(11)
        p2.font.name = 'Inter'
        p2.font.color.rgb = RGBColor(140, 180, 160) if dark else C_TEXT_MUTED

    # =========================================================================
    # SLIDE 01: COVER
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1, C_DARK_BG)
    
    # Decorative accent card
    add_shape_card(slide1, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.8), C_PINE, C_EMERALD)
    
    tb = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CANOPYAI"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    p = tf.add_paragraph()
    p.text = "Satellite-Based Urban Tree Equity Auditor & AI Plantation Prioritization Engine"
    p.font.size = Pt(22)
    p.font.color.rgb = C_EMERALD_ACCENT
    p.space_before = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "📍 Delhi, India  ·  250 Municipal Wards Analyzed"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 225, 215)
    p.space_before = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "ASEP Group 4  |  VIT Pune"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.space_before = Pt(40)

    add_footer(slide1, "The question is not whether a city needs more trees. The question is where they will create the most impact.", 1, dark=True)

    # =========================================================================
    # SLIDE 02: THE PROBLEM
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2, C_CANVAS)
    add_header(slide2, "URBAN GREEN COVER IS NOT DISTRIBUTED EQUALLY", "Delhi faces a compounding environmental crisis: severe canopy deficit in dense, vulnerable wards.")

    # 3 Stat Cards
    stats = [
        ("8.77%", "Canopy Cover Detected", "Critical tree gap in analyzed region", C_VERY_HIGH),
        ("53.26%", "Bare / Open Surface", "Unvegetated area ready for intervention", C_HIGH),
        ("21.21%", "Built-up Urban Surface", "Concrete & infrastructure heat traps", C_TEXT_MUTED)
    ]
    for i, (val, label, desc, col) in enumerate(stats):
        left = Inches(0.8 + i * 3.95)
        add_shape_card(slide2, left, Inches(1.8), Inches(3.7), Inches(3.2), C_WHITE, C_BORDER)
        tb = slide2.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.3), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = col
        
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_TEXT
        p.space_before = Pt(10)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT_MUTED
        p.space_before = Pt(6)

    # Callout Banner
    add_shape_card(slide2, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.2), C_DARK_BG)
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(5.35), Inches(11.333), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = '💡 KEY INSIGHT: "Knowing HOW MANY TREES A CITY HAS is not enough."'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    add_footer(slide2, "To decide where the next tree should go, city-wide averages are too coarse.", 2)

    # =========================================================================
    # SLIDE 03: WHY WARD-LEVEL ANALYSIS FAILS
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3, C_CANVAS)
    add_header(slide3, "THE AVERAGE CAN HIDE THE OUTLIER", "Comparing ward-level averaging against pixel-level granularity.")

    # Left: Ward View
    add_shape_card(slide3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), C_WHITE, C_BORDER)
    tb = slide3.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏛️ WARD-LEVEL VIEW"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_DARK_BG
    
    p = tf.add_paragraph()
    p.text = "Single Ward Average: 24% Canopy"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_HIGH
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Treats the entire ward as uniform\n• Hides severe zero-canopy pockets\n• Misallocates municipal planting budget\n• Cannot pinpoint exact physical locations"
    p.font.size = Pt(14)
    p.font.color.rgb = C_TEXT_MUTED
    p.space_before = Pt(12)

    # Right: Pixel View
    add_shape_card(slide3, Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.6), C_PINE)
    tb = slide3.shapes.add_textbox(Inches(7.233), Inches(2.0), Inches(5.0), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🛰️ PIXEL-LEVEL VIEW (CanopyAI)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    p = tf.add_paragraph()
    p.text = "Precise Spatial Heterogeneity"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD_ACCENT
    p.space_before = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Identifies 0% canopy sub-zones instantly\n• Separates built-up structures from plantable land\n• Overlays satellite heat islands (LST) per pixel\n• Two spots in the same ward get different treatments"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(220, 235, 225)
    p.space_before = Pt(12)

    add_footer(slide3, "So we change the unit of intelligence.", 3)

    # =========================================================================
    # SLIDE 04: CORE INNOVATION
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4, C_DARK_BG)
    add_header(slide4, "FROM WARD-LEVEL REPORTING TO PIXEL-LEVEL INTELLIGENCE", "Separating analytical precision from administrative execution.", dark=True)

    # Two Main Concept Cards
    add_shape_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.6), C_PINE, C_BLUE_ACCENT)
    tb = slide4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PIXEL"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p = tf.add_paragraph()
    p.text = "ANALYTICAL UNIT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p = tf.add_paragraph()
    p.text = "13-Band Sentinel Multispectral Pixel\nCalculates canopy, heat, and plantability at maximum spatial resolution."
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 220, 230)
    p.space_before = Pt(12)

    add_shape_card(slide4, Inches(6.933), Inches(1.8), Inches(5.6), Inches(3.6), C_EMERALD, C_EMERALD_ACCENT)
    tb = slide4.shapes.add_textbox(Inches(7.233), Inches(2.0), Inches(5.0), Inches(3.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "WARD"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD_ACCENT
    p = tf.add_paragraph()
    p.text = "ADMINISTRATIVE UNIT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p = tf.add_paragraph()
    p.text = "250 MCD Municipal Boundaries\nAggregates pixel intelligence into actionable budget, tree count, and governance reports."
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(200, 230, 215)
    p.space_before = Pt(12)

    # Core Mantra Banner
    add_shape_card(slide4, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.0), C_WHITE)
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(11.333), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = '"WE DON\'T START WITH THE WARD. WE START WITH THE PIXEL."'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_DARK_BG

    add_footer(slide4, "But a pixel alone is not a decision. We need AI to understand what that pixel represents.", 4, dark=True)

    # =========================================================================
    # SLIDE 05: HOW CANOPYAI WORKS (PIPELINE)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5, C_CANVAS)
    add_header(slide5, "FROM SATELLITE OBSERVATION TO PLANTING DECISION", "End-to-end 6-stage operational pipeline.")

    steps = [
        ("01", "Satellite Input", "13-band Sentinel GeoTIFF raster"),
        ("02", "Tiling", "256×256 patch generation"),
        ("03", "AI Engine", "13-channel SegFormer segmentation"),
        ("04", "Reconstruction", "Probability-aware tile stitching"),
        ("05", "Impact Engine", "Canopy + Heat + Plantability scoring"),
        ("06", "Decision Layer", "Priority zones & ward allocation")
    ]

    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.8 + (i % 3) * 3.95)
        top = Inches(1.8 if i < 3 else 4.2)
        add_shape_card(slide5, left, top, Inches(3.7), Inches(2.1), C_WHITE, C_BORDER)
        tb = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(3.3), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = C_DARK_BG
        
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_TEXT
        p.space_before = Pt(4)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_TEXT_MUTED
        p.space_before = Pt(4)

    add_footer(slide5, "At the center of this pipeline is the model that converts spectral information into land-cover intelligence.", 5)

    # =========================================================================
    # SLIDE 06: SYSTEM ARCHITECTURE
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6, C_CANVAS)
    add_header(slide6, "ONE PIPELINE. FOUR INTELLIGENCE LAYERS.", "Full technical software & data stack.")

    layers = [
        ("DATA LAYER", "Sentinel-2 13-Band Rasters · MCD 250 Ward GeoJSON · LST Heat Layers · NDVI", C_BLUE_ACCENT),
        ("AI SEGMENTATION LAYER", "PyTorch SegFormer (MiT-B0) · 13-Channel Adapter · Tiled Inference · Probability Stitcher", C_PINE),
        ("GEOSPATIAL ANALYTICS", "Rasterio · GeoPandas · Zonal Statistics · Plantability Index · Impact Score Engine", C_EMERALD_ACCENT),
        ("DECISION & APP LAYER", "FastAPI REST API · React 19 + Vite 8 · Leaflet GIS Viewer · Recharts · PDF Exporter", C_DARK_BG)
    ]

    for i, (title, desc, color) in enumerate(layers):
        top = Inches(1.8 + i * 1.2)
        add_shape_card(slide6, Inches(0.8), top, Inches(11.733), Inches(1.05), C_WHITE, C_BORDER)
        
        # Color bar indicator
        add_shape_card(slide6, Inches(0.8), top, Inches(0.2), Inches(1.05), color)
        
        tb = slide6.shapes.add_textbox(Inches(1.2), top + Inches(0.15), Inches(11.1), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT
        p.space_before = Pt(4)

    add_footer(slide6, "The architecture becomes powerful because the model does not stop at classification.", 6)

    # =========================================================================
    # SLIDE 07: SEGFORMER AI ENGINE
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7, C_CANVAS)
    add_header(slide7, "THE AI SEES 13 BANDS — NOT JUST RGB", "Multispectral deep semantic segmentation with SegFormer MiT-B0.")

    # Left: Channel Comparison
    add_shape_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.6), C_WHITE, C_BORDER)
    tb = slide7.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📷 Standard RGB (3 Channels)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_MUTED
    p = tf.add_paragraph()
    p.text = "Red, Green, Blue visible spectrum only. Fails to distinguish artificial turf from live tree canopy."
    p.font.size = Pt(13)
    p.font.color.rgb = C_TEXT_MUTED
    p.space_before = Pt(4)
    
    p = tf.add_paragraph()
    p.text = "🛰️ CanopyAI Sentinel-2 (13 Channels)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_DARK_BG
    p.space_before = Pt(16)
    p = tf.add_paragraph()
    p.text = "Visible (B2,B3,B4) + Red Edge (B5,B6,B7) + NIR (B8,B8A) + SWIR (B11,B12) + Coastal/Vapour (B1,B9,B10).\nDeep spectral signature extraction for high accuracy."
    p.font.size = Pt(13)
    p.font.color.rgb = C_TEXT
    p.space_before = Pt(4)

    # Right: 4 Classes
    add_shape_card(slide7, Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.6), C_PINE)
    tb = slide7.shapes.add_textbox(Inches(7.233), Inches(2.0), Inches(5.0), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏷️ 4 SEMANTIC TARGET CLASSES"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    classes_info = [
        ("0 — Bare / Open Land", "53.26% coverage · IoU 0.815"),
        ("1 — Tree Canopy", "8.77% coverage · IoU 0.626"),
        ("2 — Built-up Surface", "21.21% coverage · IoU 0.841"),
        ("3 — Cropland / Veg", "16.76% coverage · IoU 0.827")
    ]
    for cls_title, cls_desc in classes_info:
        p = tf.add_paragraph()
        p.text = f"• {cls_title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = C_EMERALD_ACCENT
        p.space_before = Pt(10)
        
        p2 = tf.add_paragraph()
        p2.text = f"   {cls_desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(210, 230, 220)

    add_footer(slide7, "Once every pixel has a semantic meaning, we can ask: which pixels deserve intervention?", 7)

    # =========================================================================
    # SLIDE 08: PIXEL IMPACT ENGINE
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8, C_CANVAS)
    add_header(slide8, "EVERY PIXEL GETS A PRIORITY", "Multi-factor spatial scoring formula (TPIS).")

    # Formula Box
    add_shape_card(slide8, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.3), C_DARK_BG)
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.333), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "S(x,y) = w_C · Deficit + w_H · Heat + w_V · NDVI Deficit + w_U · Vulnerability + w_W · Water"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD_ACCENT

    # 4 Factor Cards
    factors = [
        ("30% Low Canopy", "Canopy Deficit Index", "Where tree cover is missing", C_VERY_HIGH),
        ("25% High Temp", "LST Heat Island", "Land surface heat stress (°C)", C_HIGH),
        ("20% Veg Deficit", "NDVI Health Gap", "Vegetation health index deficit", C_MEDIUM),
        ("25% Vulnerability & Water", "UVS + Water Feasibility", "Population, slum density & water", C_LOW)
    ]
    for i, (wt, name, desc, col) in enumerate(factors):
        left = Inches(0.8 + i * 2.95)
        add_shape_card(slide8, left, Inches(3.3), Inches(2.75), Inches(3.1), C_WHITE, C_BORDER)
        tb = slide8.shapes.add_textbox(left + Inches(0.15), Inches(3.45), Inches(2.45), Inches(2.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = wt
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = col
        
        p = tf.add_paragraph()
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEXT
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = C_TEXT_MUTED
        p.space_before = Pt(6)

    add_footer(slide8, "High-scoring pixels are useful individually — but planting happens in places, not isolated pixels.", 8)

    # =========================================================================
    # SLIDE 09: PIXEL -> ZONE -> WARD
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9, C_CANVAS)
    add_header(slide9, "WE DON'T PLANT PIXELS. WE USE PIXELS TO FIND WHERE TO PLANT.", "Spatial transformation from raw data to civic implementation.")

    t_steps = [
        ("STEP 1", "Raw Raster Pixels", "Continuous pixel-level scoring across Sentinel-2 grid"),
        ("STEP 2", "High-Priority Pixels", "Thresholding pixels with Composite Score > 75"),
        ("STEP 3", "Spatial Clusters", "DBSCAN / Morphological clustering into contiguous zones"),
        ("STEP 4", "Ward Aggregation", "Zonal statistics rollup for MCD Ward 78 (Bazar Sita Ram)")
    ]

    for i, (st, title, desc) in enumerate(t_steps):
        left = Inches(0.8 + i * 2.95)
        add_shape_card(slide9, left, Inches(1.8), Inches(2.75), Inches(4.6), C_PINE if i == 3 else C_WHITE, C_BORDER)
        tb = slide9.shapes.add_textbox(left + Inches(0.15), Inches(2.0), Inches(2.45), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = st
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_EMERALD_ACCENT if i == 3 else C_TEXT_MUTED
        
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_WHITE if i == 3 else C_TEXT
        p.space_before = Pt(8)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(210, 230, 220) if i == 3 else C_TEXT_MUTED
        p.space_before = Pt(12)

    add_footer(slide9, "Now the intelligence can move from the model into a decision-maker's hands.", 9)

    # =========================================================================
    # SLIDE 10: GIS / DECISION DASHBOARD
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10, C_CANVAS)
    add_header(slide10, "FROM MODEL OUTPUT TO DECISION", "Production React 19 + Leaflet GIS Dashboard Interface.")

    # Simulated Mockup Cards
    add_shape_card(slide10, Inches(0.8), Inches(1.8), Inches(7.8), Inches(4.6), C_WHITE, C_BORDER)
    tb = slide10.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(7.2), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🗺️ INTERACTIVE MAP VIEWER (Leaflet)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_DARK_BG
    
    p = tf.add_paragraph()
    p.text = "• 250 Ward Boundary GeoJSON Chloropleth Layer\n• AI Canopy Prediction Mask Overlay (GeoTIFF)\n• Urban Heat Island (LST) Surface Overlay\n• Ward Click Inspector: Shows pop, slum density, water tier & target trees"
    p.font.size = Pt(14)
    p.font.color.rgb = C_TEXT
    p.space_before = Pt(12)

    # Right Panel: Top Ward Highlight
    add_shape_card(slide10, Inches(8.8), Inches(1.8), Inches(3.733), Inches(4.6), C_PINE)
    tb = slide10.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.333), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏆 TOP PRIORITY WARD"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD_ACCENT
    
    p = tf.add_paragraph()
    p.text = "BAZAR SITA RAM"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = "• Impact Score: 88.52 / 100\n• Trees Suggested: 14,018\n• Budget Required: ₹33.64 Lakhs\n• Cooling Benefit: -1.40 °C\n• CO₂ Absorption: 350 Tons/yr"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(210, 230, 220)
    p.space_before = Pt(12)

    add_footer(slide10, "A strong model is useful only when its output becomes understandable, explainable, and actionable.", 10)

    # =========================================================================
    # SLIDE 11: RESULTS & VALIDATION
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11, C_CANVAS)
    add_header(slide11, "HOW DO WE KNOW THE MODEL WORKS?", "Rigorous empirical performance evaluation on held-out test dataset.")

    # 4 Main Metrics Cards
    metrics = [
        ("88.71%", "Pixel Accuracy", C_EMERALD_ACCENT),
        ("0.8717", "F1 Score", C_PINE),
        ("0.7772", "Mean IoU (mIoU)", C_DARK_BG),
        ("87.33%", "Precision", C_BLUE_ACCENT)
    ]
    for i, (val, label, col) in enumerate(metrics):
        left = Inches(0.8 + i * 2.95)
        add_shape_card(slide11, left, Inches(1.8), Inches(2.75), Inches(1.8), C_WHITE, C_BORDER)
        tb = slide11.shapes.add_textbox(left + Inches(0.15), Inches(1.95), Inches(2.45), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = col
        
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEXT
        p.space_before = Pt(4)

    # Class IoU Table Card
    add_shape_card(slide11, Inches(0.8), Inches(3.9), Inches(11.733), Inches(2.5), C_WHITE, C_BORDER)
    tb = slide11.shapes.add_textbox(Inches(1.1), Inches(4.1), Inches(11.1), Inches(2.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📊 CLASS-WISE INTERSECTION OVER UNION (IoU)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_DARK_BG
    
    p = tf.add_paragraph()
    p.text = "• Built-up Surface IoU: 0.841   |   Cropland / Veg IoU: 0.827   |   Bare / Open Land IoU: 0.815   |   Tree Canopy IoU: 0.626\n• Test Set Size: Held-out Sentinel-2 tiles over Delhi Municipal Region\n• Device Performance: CUDA accelerated GPU inference (< 42s full-city execution)"
    p.font.size = Pt(14)
    p.font.color.rgb = C_TEXT_MUTED
    p.space_before = Pt(8)

    add_footer(slide11, "Model accuracy validates the AI layer. Business viability determines whether the intelligence can reach the city.", 11)

    # =========================================================================
    # SLIDE 12: BUSINESS MODEL
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_bg(slide12, C_CANVAS)
    add_header(slide12, "WE SELL BETTER PLANTING DECISIONS — NOT JUST A MAP", "Commercial B2G SaaS & ClimateTech revenue architecture.")

    tiers = [
        ("LAYER 1 — CORE (B2G SaaS)", "Municipal Corporations, Smart City SPVs, Urban Local Bodies\nProduct: AI Dashboard, Priority Zones, Ward Reports, Satellite Updates\nRevenue: Annual Subscription (₹8L – ₹75L / city / year)\n[Indicative / pilot-dependent]", C_PINE),
        ("LAYER 2 — CONSULTING & ANALYTICS", "State Urban Development Departments, Infra Developers, CSR Funds\nProduct: Custom City Assessments, Green Compliance Reports, Heat Risk Audits\nRevenue: Project-based (₹5L – ₹50L / project)\n[Illustrative]", C_EMERALD),
        ("LAYER 3 — FUTURE (CARBON MRV)", "Carbon Project Developers, Climate Funds\nProduct: Remote-sensing Plantation MRV (Monitoring, Reporting & Verification)\nRevenue: Verification Data Feeds & Impact Certification Partnerships\n[Future carbon-MRV opportunity — requires independent verification]", C_DARK_BG)
    ]

    for i, (title, desc, color) in enumerate(tiers):
        top = Inches(1.8 + i * 1.6)
        add_shape_card(slide12, Inches(0.8), top, Inches(11.733), Inches(1.4), C_WHITE, C_BORDER)
        add_shape_card(slide12, Inches(0.8), top, Inches(0.2), Inches(1.4), color)
        
        tb = slide12.shapes.add_textbox(Inches(1.2), top + Inches(0.15), Inches(11.1), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT
        p.space_before = Pt(4)

    add_footer(slide12, "The business starts with a city problem, monetizes through recurring software, and expands as trust grows.", 12)

    # =========================================================================
    # SLIDE 13: GO-TO-MARKET & SCALE
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_bg(slide13, C_CANVAS)
    add_header(slide13, "LAND WITH ONE CITY. EXPAND THROUGH THE SYSTEM.", "4-Phase expansion strategy from Delhi pilot to nationwide rollout.")

    phases = [
        ("PHASE 1: PROVE", "Delhi MCD Pilot\nValidate SegFormer model & ward recommendations with municipal authorities.", C_PINE),
        ("PHASE 2: EXPAND", "Top Indian Metros\nDeploy to Mumbai, Bengaluru, Chennai & Hyderabad urban local bodies.", C_EMERALD),
        ("PHASE 3: SCALE", "State Contracts\nState-level Urban Development & Forest Dept procurement via GeM / Tenders.", C_DARK_BG),
        ("PHASE 4: ADJACENT", "Enterprise / International\nCSR Green Belts, Highway Corridors, Heat Vulnerability & Global Climate Programs.", C_BLUE_ACCENT)
    ]

    for i, (p_title, p_desc, col) in enumerate(phases):
        left = Inches(0.8 + i * 2.95)
        add_shape_card(slide13, left, Inches(1.8), Inches(2.75), Inches(4.6), C_WHITE, C_BORDER)
        add_shape_card(slide13, left, Inches(1.8), Inches(2.75), Inches(0.6), col)
        
        # Header text inside bar
        tb = slide13.shapes.add_textbox(left + Inches(0.1), Inches(1.9), Inches(2.55), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        
        # Body text
        tb2 = slide13.shapes.add_textbox(left + Inches(0.15), Inches(2.6), Inches(2.45), Inches(3.6))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = p_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = C_TEXT

    add_footer(slide13, "Once a city uses CanopyAI for one plantation cycle, the next opportunity is the next update, ward, and city.", 13)

    # =========================================================================
    # SLIDE 14: CLOSING / DEFENSIBILITY
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_bg(slide14, C_DARK_BG)
    
    tb = slide14.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.733), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "THE MORE CITIES WE ANALYZE, THE STRONGER THE SYSTEM BECOMES"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    
    p = tf.add_paragraph()
    p.text = "🔄 Defensibility Flywheel: More Cities → More Satellite Data → Refined Regional AI Models → Higher Decision Trust"
    p.font.size = Pt(16)
    p.font.color.rgb = C_EMERALD_ACCENT
    p.space_before = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "SATELLITE  +  AI  +  GIS  +  PIXEL INTELLIGENCE  =  ACTIONABLE URBAN TREE EQUITY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_BLUE_ACCENT
    p.space_before = Pt(36)
    
    p = tf.add_paragraph()
    p.text = '"Wards tell us WHO governs an area. Pixels tell us WHERE the intervention matters."'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = C_WHITE
    p.space_before = Pt(30)
    
    p = tf.add_paragraph()
    p.text = "CANOPYAI  ·  Delhi, India  ·  ASEP Group 4 (VIT Pune)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 210, 195)
    p.space_before = Pt(36)

    add_footer(slide14, "CanopyAI turns satellite observations into planting decisions.", 14, dark=True)

    out_path = r"c:\Users\shrav\Documents\ASEP GRP4 PRO\CanopyAI\CanopyAI_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == "__main__":
    build_presentation()
